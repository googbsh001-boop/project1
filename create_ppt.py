import os
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import math

# 한국어 폰트 설정
plt.rc('font', family='Malgun Gothic')
plt.rcParams['axes.unicode_minus'] = False

def add_p(tf, text, level=0, bold=False, size=12):
    if len(tf.paragraphs) == 1 and tf.text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.level = level
    p.font.bold = bold
    p.font.size = Pt(size)
    p.font.name = 'Malgun Gothic'
    return p

def create_ppt():
    # 1. Load Data
    summary_path = r"E:\인프라수주팀\트레이닝\프로젝트1\bidding_summary.csv"
    raw_data_path = r"E:\인프라수주팀\트레이닝\프로젝트1\bidding_analysis.csv"
    
    if not os.path.exists(summary_path) or not os.path.exists(raw_data_path):
        print(f"Error: Data file not found")
        return
        
    df_summary = pd.read_csv(summary_path)
    df_summary = df_summary.dropna(subset=['mean'])
    df_summary['mean_pct'] = df_summary['mean'] * 100

    df_raw = pd.read_csv(raw_data_path)
    df_raw = df_raw.dropna(subset=['Ratio'])
    df_raw['Ratio_pct'] = df_raw['Ratio'] * 100

    winners_path = r"E:\인프라수주팀\트레이닝\프로젝트1\bidding_winners.csv"
    if os.path.exists(winners_path):
        df_winners = pd.read_csv(winners_path)
    else:
        df_winners = pd.DataFrame(columns=['File', 'WinnerCompany', 'WinnerFaction'])

    united_faction_name = '우리+무소속(주황/흰색)'
    
    df_summary['Faction'] = df_summary['Faction'].replace({'무소속(흰색)': united_faction_name, '우리(주황)': united_faction_name})
    df_raw['Faction'] = df_raw['Faction'].replace({'무소속(흰색)': united_faction_name, '우리(주황)': united_faction_name})

    def clean_filename(fname):
        name = fname.replace('입찰결과 - ', '').replace('.xlsb', '').replace('.xlsx', '')
        if '_Rev' in name:
            name = name.split('_Rev')[0]
        import re
        name = re.sub(r'^\d+\s+', '', name)
        return name
        
    df_raw['Site'] = df_raw['File'].apply(clean_filename)

    # 4:3 비율(10x7.5인치)에 맞춘 차트 사이즈 (여백 고려 약 8x4.5)
    chart_figsize = (8, 4.5)

    # 2. 세력별/업체별 막대그래프 
    plt.figure(figsize=chart_figsize)
    
    sorter = ['원숭이(하늘)', '개(그린)', united_faction_name]
    df_summary['Faction'] = pd.Categorical(df_summary['Faction'], categories=sorter, ordered=True)
    df_summary = df_summary.sort_values('Faction')

    # 업체별 낙찰 건수 계산 및 범례 이름 업데이트
    company_win_counts = df_winners['WinnerCompany'].value_counts()
    df_summary['Company'] = df_summary['Company'].apply(
        lambda x: f"{x} ({company_win_counts.get(x, 0)}건)" if pd.notna(x) else x
    )

    ax = sns.barplot(data=df_summary, x='Faction', y='mean_pct', hue='Company', palette='tab20')
    
    # 해당 세력의 누적 낙찰 갯수 계산하여 제목/범례에 추가
    winner_counts = df_winners['WinnerFaction'].replace({'무소속(흰색)': united_faction_name, '우리(주황)': united_faction_name}).value_counts()
    title_text = '세력별 / 업체별 평균 투찰률 (%) - 그룹 막대 그래프'
    sub_text = '[ 누적 낙찰 건수 | '
    for fac in sorter:
        cnt = winner_counts.get(fac, 0)
        sub_text += f"{fac}: {cnt}건, "
    sub_text = sub_text.rstrip(', ') + ' ]'
    
    plt.title(f"{title_text}\n{sub_text}", fontsize=9, fontweight='bold')
    plt.ylabel('평균 투찰률 (%)', fontsize=10)
    plt.xlabel('세력 구분 (우리/무소속 통합)', fontsize=10)
    
    # 5% (5.0) 단위 세분화
    from matplotlib.ticker import MultipleLocator
    ax.yaxis.set_major_locator(MultipleLocator(5.0))
    
    # Y축 범위를 5% 단위가 잘 보이도록 넓게 잡되, 데이터가 잘리진 않게 (예: 50~100)
    # 이미지처럼 넓은 간격을 보기 위해 최소값을 확 내림
    y_min = max(0, df_summary['mean_pct'].min() - 30)
    y_max = 100
    plt.ylim(y_min, y_max)
    
    plt.grid(axis='y', linestyle='-', alpha=0.7)
    
    plt.legend(bbox_to_anchor=(1.01, 1), loc='upper left', fontsize=9, ncol=2)
    plt.tight_layout()

    chart_path = r"E:\인프라수주팀\트레이닝\프로젝트1\faction_chart_grouped.png"
    plt.savefig(chart_path, dpi=300)
    plt.close()
    
    # 3. Create PPT (4:3 is default in python-pptx)
    prs = Presentation()

    # Slide 1: Title
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "세력별 투찰성향 분석 보고서"
    title.text_frame.paragraphs[0].font.size = Pt(15)
    title.text_frame.paragraphs[0].font.name = 'Malgun Gothic'
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = "24년도 이후 입찰결과 데이터 기반\n분석 대상: 총 27개 주요 업체 및 모든 현장 비교"
    for p in subtitle.text_frame.paragraphs:
        p.font.size = Pt(12)
        p.font.name = 'Malgun Gothic'

    # Slide 2: Overview and Summary
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "[결론] 세력별 투찰성향 요약"
    title_shape.text_frame.paragraphs[0].font.size = Pt(15)
    title_shape.text_frame.paragraphs[0].font.name = 'Malgun Gothic'
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    tf = body_shape.text_frame
    tf.text = "" 
    add_p(tf, "■ 세력 통합 업데이트: 우리(주황) 그룹과 무소속(흰색) 그룹의 성향이 유사하여 통합 분석", level=0, size=12)
    add_p(tf, "", level=0, size=12) # 한 줄 띄우기
    
    add_p(tf, "■ 원숭이(하늘) 및 개(그린) 세력: [저가, 공격적 성향]", level=0, bold=False, size=12)
    add_p(tf, "- 85%~87%대 중후반의 매우 낮은 평균 투찰률(기초대비)을 형성", level=1, size=12)
    add_p(tf, "- 대표 업체: 동부건설(85.8%), 대우건설(87.0%), 계룡건설산업(87.1%) 등", level=1, size=12)
    add_p(tf, "", level=0, size=12) # 한 줄 띄우기
    
    add_p(tf, "■ 우리+무소속 통합 세력: [보수적, 수익성 지향형]", level=0, bold=False, size=12)
    add_p(tf, "- 비교적 높은 투찰률을 보이며 무리한 출혈 경쟁 지양, 안정적 이익 확보 우선", level=1, size=12)
    add_p(tf, "- 대표 업체: 한화(93.5%), 대보건설(93.3%), 호반산업(93.0%) 등", level=1, size=12)

    # Slide 3: Chart
    blank_slide_layout = prs.slide_layouts[5]  # Title only
    slide = prs.slides.add_slide(blank_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "📊 통합 세력별 평균 투찰률(%) 분석 (그룹 막대그래프)"
    title_shape.text_frame.paragraphs[0].font.size = Pt(15)
    title_shape.text_frame.paragraphs[0].font.name = 'Malgun Gothic'
    title_shape.text_frame.paragraphs[0].font.bold = True

    # Center for 4:3 (Width 10, Height 7.5)
    pic_width = Inches(9.0)
    left = Inches(0.5)
    top = Inches(1.8)
    pic = slide.shapes.add_picture(chart_path, left, top, width=pic_width)
    
    # -----------------------------------------------------------------
    df_winners['Site'] = df_winners['File'].astype(str).apply(clean_filename)
    import re
    def get_date_val(site):
        orig_file = ""
        cw = df_winners[df_winners['Site'] == site]
        if not cw.empty:
            orig_file = cw.iloc[0]['File']
        else:
            cr = df_raw[df_raw['Site'] == site]
            if not cr.empty:
                orig_file = cr.iloc[0]['File']
        
        if orig_file:
            match = re.search(r'(?:입찰결과\s*-\s*)?(\d{2})(\d{2})(\d{2})', str(orig_file))
            if match:
                yy = int(match.group(1))
                mm = int(match.group(2))
                dd = int(match.group(3))
                year = 2000 + yy if yy < 50 else 1900 + yy
                return year * 10000 + mm * 100 + dd
        return 99999999

    unique_sites = sorted(df_raw['Site'].unique(), key=get_date_val)
    chunk_size = 4 
    total_pages = math.ceil(len(unique_sites) / chunk_size)
    
    for i in range(total_pages):
        chunk_sites = unique_sites[i*chunk_size : (i+1)*chunk_size]
        df_site_filtered = df_raw[df_raw['Site'].isin(chunk_sites)].copy()
        
        site_faction_mean = df_site_filtered.groupby(['Site', 'Faction'])['Ratio_pct'].mean().reset_index()
        site_faction_mean['Faction'] = pd.Categorical(site_faction_mean['Faction'], categories=sorter, ordered=True)
        
        plt.figure(figsize=chart_figsize)
        ax = sns.barplot(data=site_faction_mean, x='Site', y='Ratio_pct', hue='Faction', palette='Set2')
        
        plt.title(f'개별 현장별 세력 경쟁 및 투찰률 추이 ({i+1}/{total_pages})', fontsize=12, fontweight='bold')
        plt.ylabel('투찰률 (%)', fontsize=10)
        
        # X축 라벨(현장명)에 낙찰사 정보 및 입찰 연/월 추가
        new_labels = []
        import textwrap
        import re
        for site in chunk_sites:
            winner = "미상"
            orig_file_for_date = ""
            for _, wrow in df_winners.iterrows():
                if clean_filename(wrow['File']) == site:
                    winner = wrow['WinnerCompany']
                    orig_file_for_date = wrow['File']
                    break
                    
            if not orig_file_for_date:
                orig_files = df_raw[df_raw['Site'] == site]['File'].unique()
                if len(orig_files) > 0:
                    orig_file_for_date = orig_files[0]
                    
            date_str = ""
            if orig_file_for_date:
                match = re.search(r'(?:입찰결과\s*-\s*)?(\d{2})(\d{2})\d{2}', orig_file_for_date)
                if match:
                    yy = int(match.group(1))
                    mm = int(match.group(2))
                    year = 2000 + yy if yy < 50 else 1900 + yy
                    date_str = f"[{year}년 {mm}월]"

            # 사이트 이름이 너무 길면 줄바꿈, 아래 날짜와 낙찰사 추가
            wrapped_site = textwrap.fill(site, width=18)
            if date_str:
                new_labels.append(f"{wrapped_site}\n{date_str}\n[낙찰: {winner}]")
            else:
                new_labels.append(f"{wrapped_site}\n[낙찰: {winner}]")
            
        ax.set_xticks(range(len(new_labels)))
        ax.set_xticklabels(new_labels, rotation=10, ha='center', fontsize=8)
        plt.xlabel('입찰 현장명', fontsize=10)
        
        # 5% 단위
        plt.gca().yaxis.set_major_locator(MultipleLocator(5.0))
        
        # 막대그래프 위에 겹치지 않게 작은 사이즈로 표시
        for container in ax.containers:
            ax.bar_label(container, fmt='%.2f%%', padding=2, fontsize=6)
        
        y_min_site = max(0, site_faction_mean['Ratio_pct'].min() - 30)
        y_max_site = 100
        plt.ylim(y_min_site, y_max_site)
        
        plt.grid(axis='y', linestyle='-', alpha=0.7)
        plt.legend(bbox_to_anchor=(1.01, 1), loc='upper left', fontsize=9)
        plt.tight_layout()

        site_chart_path = r"E:\인프라수주팀\트레이닝\프로젝트1\site_chart_chunk_" + str(i) + ".png"
        plt.savefig(site_chart_path, dpi=300)
        plt.close()
        
        slide = prs.slides.add_slide(blank_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = f"🏗️ 전체 현장별/세력별 투찰 성향 상세 비교 ({i+1}/{total_pages})"
        title_shape.text_frame.paragraphs[0].font.size = Pt(15)
        title_shape.text_frame.paragraphs[0].font.name = 'Malgun Gothic'
        title_shape.text_frame.paragraphs[0].font.bold = True

        slide.shapes.add_picture(site_chart_path, left, top, width=pic_width)

    # Slide Last: Strategy
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "💡 현장/경쟁사 맞춤형 전략 제언"
    title_shape.text_frame.paragraphs[0].font.size = Pt(15)
    title_shape.text_frame.paragraphs[0].font.name = 'Malgun Gothic'
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    tf = body_shape.text_frame
    tf.text = ""
    add_p(tf, "■ 적자 수주 가능성 등 리스크 관리", level=0, bold=False, size=12)
    add_p(tf, "- 막연히 최저가를 추종할 경우 이익 훼손 소지가 큼", level=1, size=12)
    add_p(tf, "", level=0, size=12) # 띄어쓰기

    add_p(tf, "■ 경쟁사 맞춤형(Targeting) 입찰 전략 수립 제안", level=0, bold=False, size=12)
    add_p(tf, "", level=0, size=12) # 띄어쓰기

    add_p(tf, "1. 초공격적 성향의 업체 포진 시 (원숭이/개 세력 주도)", level=1, size=12)
    add_p(tf, "▶ 이전 시뮬레이션 모델을 활용하여 [한계 원가 기반 최대한의 공격적 밴드]로 하향 조정 투찰", level=2, size=12)
    add_p(tf, "", level=0, size=12) # 띄어쓰기

    add_p(tf, "2. 보수적 성향의 업체 포진 시 (우리+무소속 계열 주도)", level=1, size=12)
    add_p(tf, "▶ 맹목적인 가격 출혈을 멈추고 [적정 마진을 남길 수 있는 안정적 상향 투찰]로 수익성 극대화 모색", level=2, size=12)

    # Save
    output_path = r"E:\인프라수주팀\트레이닝\프로젝트1\분석결과_보고서_최종_4x3.pptx"
    try:
        prs.save(output_path)
        print(f"PPT 생성 완료: {output_path}")
    except PermissionError:
        output_path_v2 = r"E:\인프라수주팀\트레이닝\프로젝트1\분석결과_보고서_최종_4x3_v2.pptx"
        prs.save(output_path_v2)
        print(f"PPT 생성 완료: {output_path_v2}")

if __name__ == "__main__":
    create_ppt()
